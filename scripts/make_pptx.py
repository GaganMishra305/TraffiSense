#!/usr/bin/env python3
"""Generate the TraffiSense.ai judge pitch deck as a real .pptx.

Mirrors docs/deck/index.html: dark command-center theme, big stat tiles, the
leakage story, and embedded product screenshots. Output: docs/deck/TraffiSense.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
SHOTS = HERE.parent / "docs" / "screenshots"
OUT = HERE.parent / "docs" / "deck" / "TraffiSense.pptx"

# Palette
INK = RGBColor(0x0B, 0x12, 0x20)
PANEL = RGBColor(0x16, 0x23, 0x3C)
EDGE = RGBColor(0x2A, 0x3A, 0x5C)
ACCENT = RGBColor(0x22, 0xD3, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x9C, 0xAC, 0xC4)
GOOD = RGBColor(0x34, 0xD3, 0x99)
BAD = RGBColor(0xF8, 0x71, 0x71)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]
_page = [0]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = INK
    bg.line.fill.background()
    bg.shadow.inherit = False
    # top accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, Pt(6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False
    _page[0] += 1
    _footer(s)
    return s


def _footer(s):
    tb = s.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(6), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "TraffiSense.ai  ·  Predictive Traffic Command"
    r.font.size = Pt(10)
    r.font.color.rgb = MUTED
    n = s.shapes.add_textbox(Inches(11.9), Inches(7.05), Inches(1.2), Inches(0.4))
    pp = n.text_frame.paragraphs[0]
    pp.alignment = PP_ALIGN.RIGHT
    rr = pp.add_run()
    rr.text = f"{_page[0]} / 11"
    rr.font.size = Pt(10)
    rr.font.color.rgb = MUTED


def text(s, x, y, w, h, runs, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, spacing=1.0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [(runs, color, bold)]
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    for t, c, b in runs:
        r = p.add_run()
        r.text = t
        r.font.size = Pt(size)
        r.font.bold = b
        r.font.color.rgb = c
        r.font.name = "Helvetica Neue"
    return tb


def kicker(s, label):
    text(s, Inches(0.6), Inches(0.5), Inches(8), Inches(0.4),
         label.upper(), size=13, color=ACCENT, bold=True)


def title(s, t, y=Inches(0.95)):
    text(s, Inches(0.6), y, Inches(12.1), Inches(1.0), t, size=40, color=WHITE, bold=True)


def card(s, x, y, w, h, fill=PANEL, edge=EDGE):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.line.color.rgb = edge
    c.line.width = Pt(1)
    c.shadow.inherit = False
    return c


def bullets(s, x, y, w, items, size=18, gap=0.62):
    for i, (txt, col) in enumerate(items):
        tb = s.shapes.add_textbox(x, y + Inches(i * gap), w, Inches(gap))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        d = p.add_run()
        d.text = "▸  "
        d.font.size = Pt(size)
        d.font.color.rgb = ACCENT
        d.font.bold = True
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.color.rgb = col
        r.font.name = "Helvetica Neue"


def stat_tile(s, x, y, w, h, value, label, sub):
    card(s, x, y, w, h)
    text(s, x, y + Inches(0.25), w, Inches(0.9), value, size=44, color=ACCENT,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, x, y + Inches(1.15), w, Inches(0.4), label, size=15, color=WHITE,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, x, y + Inches(1.55), w, Inches(0.4), sub, size=11, color=MUTED,
         align=PP_ALIGN.CENTER)


def image_card(s, img, x, y, w):
    pic = s.shapes.add_picture(str(img), x, y, width=w)
    pic.line.color.rgb = EDGE
    pic.line.width = Pt(1.5)
    return pic


# ---------------------------------------------------------------------------
# Slide 1 - Title
s = slide()
icon = card(s, Inches(6.07), Inches(1.7), Inches(1.2), Inches(1.2), fill=PANEL, edge=ACCENT)
text(s, Inches(6.07), Inches(1.95), Inches(1.2), Inches(0.7), "T", size=44,
     color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
text(s, Inches(0.6), Inches(3.1), Inches(12.1), Inches(1.2),
     [("TraffiSense", WHITE, True), (".ai", ACCENT, True)], size=66, bold=True,
     align=PP_ALIGN.CENTER)
text(s, Inches(0.6), Inches(4.4), Inches(12.1), Inches(0.6),
     [("Predict. Simulate. Deploy.  ", WHITE, True),
      ("AI traffic command for city events", ACCENT, True)],
     size=22, align=PP_ALIGN.CENTER)
text(s, Inches(0.6), Inches(5.2), Inches(12.1), Inches(0.5),
     "Built on 8,057 real Bengaluru ASTraM incidents  ·  MapMyIndia mapping",
     size=15, color=MUTED, align=PP_ALIGN.CENTER)

# Slide 2 - Problem
s = slide()
kicker(s, "The Problem")
title(s, "Traffic management today is reactive")
bullets(s, Inches(0.6), Inches(2.0), Inches(12), [
    ("Officers watch CCTV and respond only AFTER the jam forms.", WHITE),
    ("No reliable pre-event forecasting — impact is estimated by gut feel.", MUTED),
    ("Plans (diversions, barricades, signals) deployed without network insight.", MUTED),
    ("Manpower placement is manual guesswork; emergency response slows.", MUTED),
    ("No fast 'what-if' planning for attendance, weather or road closures.", MUTED),
])
qc = card(s, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.1), fill=PANEL, edge=ACCENT)
text(s, Inches(0.9), Inches(5.62), Inches(11.6), Inches(0.9),
     [("Core question:  ", ACCENT, True),
      ("How can historical + real-time data forecast event traffic impact and "
       "recommend optimal manpower, barricading and diversion plans?", WHITE, False)],
     size=16, anchor=MSO_ANCHOR.MIDDLE)

# Slide 3 - Solution
s = slide()
kicker(s, "Our Solution")
title(s, "Reactive becomes proactive")
text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(0.8),
     "One click on an event forecasts impact, simulates strategies, optimizes "
     "resources, and plans emergency corridors — all explainable.",
     size=18, color=MUTED)
caps = [("Forecast", "Congestion before it forms"),
        ("Simulate", "Test every plan safely"),
        ("Optimize", "Right resources, right place"),
        ("Explain", "Every number has a reason")]
for i, (h, d) in enumerate(caps):
    x = Inches(0.6 + i * 3.05)
    card(s, x, Inches(3.1), Inches(2.85), Inches(1.7))
    text(s, x, Inches(3.35), Inches(2.85), Inches(0.5), h, size=20, color=ACCENT,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, x, Inches(3.95), Inches(2.85), Inches(0.8), d, size=14, color=WHITE,
         align=PP_ALIGN.CENTER)
text(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(0.6),
     "\u201cAnticipate congestion, test every plan, deploy with confidence.\u201d",
     size=20, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# Slide 4 - Five features
s = slide()
kicker(s, "Capabilities")
title(s, "Five deep features")
feats = [
    ("1  Event Intelligence", "Real venues, attendance, time, weather."),
    ("2  AI Impact Prediction", "Gravity demand x history x ML closure-risk."),
    ("3  Live Command Map", "Heatmap, congestion, hotspots, diversions, emergency. MapMyIndia + OSM fallback."),
    ("4  Resource Optimization", "Officers/barricades/marshals/tow trucks by need, each with rationale."),
    ("5  Scenarios + NL Assistant", "Compare 5 strategies; ask live what-if questions."),
]
for i, (h, d) in enumerate(feats):
    col = i % 3
    row = i // 3
    x = Inches(0.6 + col * 4.07)
    y = Inches(2.1 + row * 2.35)
    card(s, x, y, Inches(3.85), Inches(2.1))
    text(s, x + Inches(0.2), y + Inches(0.2), Inches(3.5), Inches(0.5), h,
         size=17, color=ACCENT, bold=True)
    text(s, x + Inches(0.2), y + Inches(0.85), Inches(3.5), Inches(1.1), d,
         size=13.5, color=WHITE)

# Slide 5 - Forecast screenshot
s = slide()
kicker(s, "The Forecast")
title(s, "One click, fully explained")
image_card(s, SHOTS / "demo-02-forecast.png", Inches(2.2), Inches(2.0), Inches(8.9))
text(s, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.6),
     "e.g. 1,300 vehicles onto CBD 2 -> 61% saturation, 19 min delay, 40% closure risk (ML).",
     size=14, color=MUTED, align=PP_ALIGN.CENTER)

# Slide 6 - ML stat tiles
s = slide()
kicker(s, "The Machine Learning")
title(s, "An honest, imbalanced operational target")
text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(0.8),
     [("RandomForest", WHITE, True),
      (" (200 trees, depth 12, class-balanced) predicting ", MUTED, False),
      ("road-closure requirement", WHITE, True),
      (" — a real imbalanced problem at a 7.4% base rate. Metrics real & reproducible.",
       MUTED, False)], size=16)
tiles = [("0.798", "Holdout ROC-AUC", "generalization"),
         ("0.321", "Avg Precision", "4.3x lift over base"),
         ("0.756", "5-fold CV AUC", "+/- 0.022 stable"),
         ("0.92", "Accuracy @ best-F1", "operating point")]
for i, (v, l, sub) in enumerate(tiles):
    stat_tile(s, Inches(0.6 + i * 3.05), Inches(3.3), Inches(2.85), Inches(2.1), v, l, sub)

# Slide 7 - Leakage + feature importances
s = slide()
kicker(s, "ML Maturity")
title(s, [("We caught ", WHITE, True), ("data leakage", ACCENT, True)])
text(s, Inches(0.6), Inches(2.0), Inches(6), Inches(0.8),
     "We rejected two trap targets. Honest metrics = what the model actually generalizes.",
     size=15, color=MUTED)
lk = [("priority (High/Low)", "Deterministic corridor label — faked AUC 0.998 by memorizing the corridor."),
      ("closure_flag as feature", "Literally the target in disguise.")]
for i, (h, d) in enumerate(lk):
    y = Inches(2.95 + i * 1.35)
    card(s, Inches(0.6), y, Inches(5.6), Inches(1.2), edge=BAD)
    text(s, Inches(0.85), y + Inches(0.15), Inches(5.1), Inches(0.4), h, size=15,
         color=BAD, bold=True)
    text(s, Inches(0.85), y + Inches(0.6), Inches(5.1), Inches(0.5), d, size=12.5,
         color=WHITE)
# feature importance bars
text(s, Inches(6.7), Inches(2.0), Inches(6), Inches(0.4), "Top feature importances",
     size=15, color=WHITE, bold=True)
fi = [("corridor_volume", 0.172), ("hour (sin)", 0.129), ("hour (cos)", 0.126),
      ("cause = breakdown", 0.113), ("cause = tree_fall", 0.112), ("weekday", 0.105)]
maxv = 0.172
for i, (name, v) in enumerate(fi):
    y = Inches(2.6 + i * 0.62)
    text(s, Inches(6.7), y, Inches(2.0), Inches(0.4), name, size=11.5, color=MUTED,
         anchor=MSO_ANCHOR.MIDDLE)
    bw = Inches(3.6 * v / maxv)
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.75), y + Inches(0.05),
                             bw, Inches(0.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False
    text(s, Inches(8.75) + bw + Emu(40000), y, Inches(0.8), Inches(0.4), f"{v:.2f}",
         size=10.5, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

# Slide 8 - Fusion / explainability
s = slide()
kicker(s, "Explainable by construction")
title(s, "From model to decision")
card(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(1.0), fill=PANEL, edge=ACCENT)
text(s, Inches(0.85), Inches(2.12), Inches(11.6), Inches(0.8),
     [("congestion = ", MUTED, False),
      ("0.42·time-of-day  +  0.43·event demand (gravity)  +  0.15·historical risk", WHITE, True),
      ("   × weather", MUTED, False)], size=16, anchor=MSO_ANCHOR.MIDDLE)
steps = [("ML closure-risk", "RandomForest propensity"),
         ("Digital-twin sim", "travel, queue, fuel, CO2, emergency"),
         ("Optimizer", "resources by marginal need"),
         ("NL Assistant", "live what-if answers")]
for i, (h, d) in enumerate(steps):
    x = Inches(0.6 + i * 3.05)
    card(s, x, Inches(3.5), Inches(2.85), Inches(1.7))
    text(s, x, Inches(3.7), Inches(2.85), Inches(0.5), h, size=16, color=ACCENT,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, x, Inches(4.3), Inches(2.85), Inches(0.8), d, size=12.5, color=WHITE,
         align=PP_ALIGN.CENTER)
text(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(0.5),
     "Every output traces back to evidence.", size=16, color=MUTED, align=PP_ALIGN.CENTER)

# Slide 9 - Montage
s = slide()
kicker(s, "In Action")
title(s, "Optimize · Compare · Ask")
mont = ["demo-03-resources.png", "demo-04-scenarios.png", "demo-05-assistant.png"]
for i, img in enumerate(mont):
    image_card(s, SHOTS / img, Inches(0.6 + i * 4.15), Inches(2.2), Inches(3.9))
text(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(0.8),
     "Full optimization -> 30% faster travel. Ask \u201cwhat if heavy rain?\u201d -> "
     "worst corridor 61% to 100%, delay 15 to 59 min.", size=15, color=MUTED,
     align=PP_ALIGN.CENTER)

# Slide 10 - Tech & Data
s = slide()
kicker(s, "Under the Hood")
title(s, "Tech & Data")
text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(0.4), "Stack", size=16,
     color=ACCENT, bold=True)
chips = ["FastAPI", "scikit-learn", "pandas", "Leaflet + MapMyIndia (Mappls)",
         "OSM fallback", "Chart.js", "Tailwind", "Self-contained NL assistant"]
x, y = 0.6, 2.5
for c in chips:
    w = Inches(0.35 + len(c) * 0.105)
    ch = card(s, Inches(x), Inches(y), w, Inches(0.5), fill=PANEL)
    text(s, Inches(x), Inches(y) + Inches(0.06), w, Inches(0.4), c, size=12.5,
         color=WHITE, align=PP_ALIGN.CENTER)
    x += 0.35 + len(c) * 0.105 + 0.25
    if x > 11.5:
        x = 0.6
        y += 0.7
data = [("8,057", "cleaned incidents (Nov 23 - Apr 24)"),
        ("21", "arterial corridors"),
        ("223", "chronic junction hotspots")]
for i, (v, l) in enumerate(data):
    stat_tile(s, Inches(0.6 + i * 4.07), Inches(4.4), Inches(3.85), Inches(1.7), v, l, "")
text(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.5),
     "Includes a 3-min narrated demo video with accessible burned-in captions.",
     size=13, color=MUTED, align=PP_ALIGN.CENTER)

# Slide 11 - Why it wins
s = slide()
kicker(s, "Why It Wins")
title(s, "Reactive traffic management, made proactive")
wins = [("Proactive, not reactive", "Forecast and plan before the event."),
        ("Explainable AI", "Every recommendation has a rationale operators trust."),
        ("Honest, leakage-free ML", "AUC 0.80 that actually generalizes."),
        ("Live & real", "Real ASTraM data, real maps, end-to-end.")]
for i, (h, d) in enumerate(wins):
    col = i % 2
    row = i // 2
    x = Inches(0.6 + col * 6.15)
    y = Inches(2.2 + row * 1.5)
    card(s, x, y, Inches(5.95), Inches(1.25))
    text(s, x + Inches(0.25), y + Inches(0.18), Inches(5.5), Inches(0.4), h,
         size=17, color=ACCENT, bold=True)
    text(s, x + Inches(0.25), y + Inches(0.65), Inches(5.5), Inches(0.5), d,
         size=13.5, color=WHITE)
text(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(0.7),
     "TraffiSense.ai — anticipate, test every plan, deploy with confidence.",
     size=20, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

prs.save(str(OUT))
print(f"Saved {OUT} ({len(prs.slides._sldIdLst)} slides)")
